from flask import Flask, request, send_file, render_template, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import google.generativeai as genai
from dotenv import load_dotenv
import os
import io

app = Flask(__name__)

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel('gemini-2.0-flash')

# Slide creation functions (unchanged from your code)
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def generate_slide_content(section, topic, model):
    try:
        prompt = f"Write a brief paragraph (3-4 sentences) about '{section}' in the context of {topic}."
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating content for {section}: {e}")
        return f"Content for {section} could not be generated. Please add details manually."

def generate_sections(topic, model):
    try:
        prompt = f"Generate 12-13 key sections/subtopics for a presentation about {topic}. Return as a comma-separated list."
        response = model.generate_content(prompt)
        return [s.strip() for s in response.text.split(',')]
    except Exception as e:
        print(f"Error generating sections: {e}")
        return [
            "Introduction", "Key Concepts", "Current Trends",
            "Applications", "Case Studies", "Benefits",
            "Challenges", "Future Directions", "Conclusion"
        ]

def add_welcome_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, RGBColor(13, 71, 161))
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    add_decorative_elements(slide)

def add_content_page(prs, sections, topic):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
    set_slide_background(slide, RGBColor(240, 240, 240))
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"Presentation Outline: {topic}"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(13, 71, 161)
    p.alignment = PP_ALIGN.LEFT
    for i, section in enumerate(sections):
        row = i // 3
        col = i % 3
        left_pos = Inches(0.5 + col * 3)
        top_pos = Inches(1.8 + row * 1.8)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_pos, Inches(2.8), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(13, 71, 161)
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        p = tf.add_paragraph()
        p.text = f"{i+1}. {section}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
    add_watermark(slide)

def add_content_slide(prs, title_text, content_text, topic):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
    set_slide_background(slide, RGBColor(255, 255, 255))
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(13, 71, 161)
    p.alignment = PP_ALIGN.LEFT
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, Inches(1.3), width, Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 152, 0)
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    content_box.line.color.rgb = RGBColor(200, 200, 200)
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = content_text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(33, 33, 33)
    add_watermark(slide)

def add_thank_you_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
    set_slide_background(slide, RGBColor(13, 71, 161))
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You"
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    add_decorative_elements(slide)

def add_decorative_elements(slide):
    for i in range(3):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5 + i*3), Inches(6.5), Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 152, 0)
        shape.line.color.rgb = RGBColor(255, 152, 0)

def add_watermark(slide):
    left = Inches(8)
    top = Inches(6.8)
    width = Inches(2)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Generated by SuuSri AI"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.RIGHT

def create_presentation(title, topic):
    prs = Presentation()
    sections = generate_sections(topic, model)
    add_welcome_slide(prs, title)
    add_content_page(prs, sections, topic)
    for section in sections:
        content = generate_slide_content(section, topic, model)
        add_content_slide(prs, section, content, topic)
    add_thank_you_slide(prs)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Flask routes
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate_ppt():
    data = request.get_json()
    title = data.get('title')
    topic = data.get('topic')

    if not title or not topic:
        return jsonify({'error': 'Title and topic are required'}), 400

    try:
        buffer = create_presentation(title, topic)
        return send_file(
            buffer,
            as_attachment=True,
            download_name=f"{title.replace(' ', '_')}_Presentation.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5000)